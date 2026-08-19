"""One-off generator for PE_Dashboard_Overview.pptx (2-slide exec overview + agenda).
Run: py -3.14 _make_pptx_overview.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0B, 0x14, 0x22)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xF5, 0xF7, 0xFA)
MUTED = RGBColor(0xA8, 0xB3, 0xC2)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
GREEN = RGBColor(0x34, 0xD3, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.shapes.add_shape(1, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=WHITE, gap=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (head, sub) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"{head}"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = CYAN
        run.font.name = "Segoe UI"
        if sub:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(gap + 6)
            r2 = p2.add_run()
            r2.text = sub
            r2.font.size = Pt(size - 3)
            r2.font.color.rgb = MUTED
            r2.font.name = "Segoe UI"
    return box


# ---------------------------------------------------------------- Slide 1 --
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)
add_text(s1, Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.7),
          "PE Audit Dashboard", 40, WHITE, bold=True)
add_text(s1, Inches(0.72), Inches(1.28), Inches(11.9), Inches(0.5),
          "One tool to check every customer's batch jobs, servers, and contract — in minutes, not days", 18, CYAN)

# accent rule
rule = s1.shapes.add_shape(1, Inches(0.72), Inches(1.85), Inches(2.6), Pt(3))
rule.fill.solid(); rule.fill.fore_color.rgb = AMBER; rule.line.fill.background(); rule.shadow.inherit = False

agenda_items = [
    ("1. What is it?", "A dashboard that checks a customer's batch jobs, servers, and contract terms, and tells us what's wrong — used for all 250–300 customers, no more manual Excel work."),
    ("2. How does it work?", "Upload the files → tool does all the math → shows problems found → writes the report. One flow, same steps every time."),
    ("3. What does it check?", "6 things: Batch jobs, Servers, SLA times, Contract (SOW), Benchmarks, Issues list — all in one place, not six separate tools."),
    ("4. Why can we trust the numbers?", "Every number shows where it came from. Nothing is guessed or made up. If a job is left out, it tells you why."),
    ("5. What's next?", "Live walkthrough with a real customer's files, then questions."),
]
add_bullets(s1, Inches(0.72), Inches(2.25), Inches(11.9), Inches(4.8), agenda_items, size=18)

add_text(s1, Inches(0.72), Inches(7.02), Inches(11.9), Inches(0.4),
          "Performance Engineering  ·  Internal Overview", 11, MUTED)

# ---------------------------------------------------------------- Slide 2 --
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_text(s2, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.65),
          "Key Things To Know", 34, WHITE, bold=True)
rule2 = s2.shapes.add_shape(1, Inches(0.72), Inches(1.15), Inches(2.6), Pt(3))
rule2.fill.solid(); rule2.fill.fore_color.rgb = AMBER; rule2.line.fill.background(); rule2.shadow.inherit = False

col_w = Inches(5.75)
left_items = [
    ("What it does", "Checks if a customer's batch jobs finished on time, servers are healthy, and actual usage matches the signed contract (SOW) — one dashboard per customer."),
    ("One number, everywhere", "Every screen shows the same numbers, calculated once. No panel does its own separate math, so nothing ever disagrees with itself."),
    ("SLA time comes from the right place", "Order: customer's own SLA file first, then the signed contract, then our standard default. Every value says exactly which one it used."),
    ("Fully explainable", "Any job left out of the results shows a plain-English reason — auto-detected or typed in by the reviewer. Nothing is hidden or unexplained."),
]
right_items = [
    ("6 uploads, 1 result", "Batch report, Server report, SLA file, Contract (SOW), Benchmark, Issues list — upload all 6, get one combined set of findings, not six separate reports."),
    ("~90 built-in checks", "14 groups of rules automatically flag late jobs, unhealthy servers, contract mismatches, repeat problems, and missing evidence — so nothing gets missed by hand."),
    ("No guessing, no fake data", "Every number is pulled from the customer's real uploaded files. Nothing is hardcoded or made up, so the same file always gives the same answer."),
    ("What you get", "One dashboard, a plain-English summary for leadership, and a report you can export and send — done in minutes instead of a full day of manual Excel work."),
]
add_bullets(s2, Inches(0.72), Inches(1.55), col_w, Inches(5.3), left_items, size=14)
add_bullets(s2, Inches(6.85), Inches(1.55), col_w, Inches(5.3), right_items, size=14)

add_text(s2, Inches(0.72), Inches(7.02), Inches(11.9), Inches(0.4),
          "Performance Engineering  ·  Internal Overview", 11, MUTED)

out_path = "PE_Dashboard_Overview.pptx"
try:
    prs.save(out_path)
except PermissionError:
    out_path = "PE_Dashboard_Overview_new.pptx"
    prs.save(out_path)
print(f"Saved {out_path}")
